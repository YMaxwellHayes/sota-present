#!/usr/bin/env python3
"""
sota-present 代码路径压测 / stress test.

Exercises the deterministic parts of the skill at scale and reports pass/fail:
  1. Catalog integrity      — JSON parses + cross-references resolve
  2. Palette SVG validation — build a canonical SVG per palette, run the shipped validator
  3. Validator correctness  — known-good passes, known-bad SVGs are rejected
  4. Template render        — every gallery template.html renders to a non-blank PNG (headless Chrome)
  5. Scripts smoke          — preflight / render / serve behave on good & bad input

Usage: python3 scripts/stress-test.py [--no-render]
Exit code 0 if all hard checks pass, 1 otherwise.
"""
import json, os, re, subprocess, sys, time, tempfile, glob

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
os.chdir(ROOT)
NO_RENDER = "--no-render" in sys.argv

CHROME_CANDIDATES = [
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    "/Applications/Chromium.app/Contents/MacOS/Chromium",
]
CHROME = next((c for c in CHROME_CANDIDATES if os.path.exists(c)), None)

results = []  # (section, name, ok, detail)
def rec(section, name, ok, detail=""):
    results.append((section, name, ok, detail))

HEX = re.compile(r"^#[0-9A-Fa-f]{6}$")

# ---------- 1. CATALOG INTEGRITY ----------
def check_catalog():
    try:
        styles = json.load(open("catalog/styles.json"))["styles"]
        slides = json.load(open("catalog/slides-index.json"))["slides"]
        palettes = json.load(open("catalog/whiteboard-index.json"))["palettes"]
    except Exception as e:
        rec("catalog", "parse JSON", False, str(e)); return
    rec("catalog", "parse 3 indexes", True, f"styles={len(styles)} slides={len(slides)} palettes={len(palettes)}")

    # count sanity — guard against silent shrinkage (deleting items must trip a failure, not pass quietly)
    EXPECT = {"styles": 69, "slides": 46, "palettes": 35}
    got = {"styles": len(styles), "slides": len(slides), "palettes": len(palettes)}
    rec("catalog", "counts match expected", got == EXPECT, f"expected {EXPECT}, got {got}")
    nvd = sum(1 for s in styles if s.get("verified_dual"))
    rec("catalog", "verified_dual count == 12", nvd == 12, f"got {nvd}")

    palette_ids = {p["id"] for p in palettes}

    # slides-index template_path must exist; gallery ones need template.html
    bad = []
    for s in slides:
        tp = s.get("template_path", "")
        if s.get("source") == "gallery":
            if not os.path.isdir(tp) or not os.path.isfile(os.path.join(tp, "template.html")):
                bad.append(f"{s['id']}: gallery missing dir/template.html ({tp})")
        else:  # preset: defined inside a single source file
            if not tp or not os.path.isfile(tp):
                bad.append(f"{s['id']}: preset source file missing ({tp})")
    rec("catalog", "slides-index paths resolve", not bad, "; ".join(bad[:5]) or f"{len(slides)} ok")

    # palette color roles + valid hex
    bad = []
    for p in palettes:
        cols = p.get("colors", {})
        for role in ("primary", "background", "text"):
            v = cols.get(role)
            if not v or not HEX.match(str(v)):
                bad.append(f"{p['id']}.{role}={v}")
    rec("catalog", "palette colors valid hex", not bad, "; ".join(bad[:5]) or f"{len(palettes)} ok")

    # styles cross-refs
    dangling_tpl, dangling_pal = [], []
    for st in styles:
        cap = st.get("capabilities", {})
        ref = (cap.get("slides") or {}).get("template_ref")
        if ref and not os.path.isdir(os.path.join("templates/slides", ref)):
            dangling_tpl.append(f"{st['id']}->{ref}")
        pref = (cap.get("whiteboard") or {}).get("palette_ref")
        if pref and pref not in palette_ids:
            dangling_pal.append(f"{st['id']}->{pref}")
    rec("catalog", "style.template_ref resolve", not dangling_tpl, "; ".join(dangling_tpl[:5]) or "ok")
    rec("catalog", "style.palette_ref resolve", not dangling_pal, "; ".join(dangling_pal[:5]) or "ok")

    # verified_dual must have both sides
    bad = []
    vd = [s for s in styles if s.get("verified_dual")]
    for st in vd:
        cap = st.get("capabilities", {})
        sl = (cap.get("slides") or {})
        has_slide = bool(sl.get("template_ref") or sl.get("preset_ref"))
        has_pal = bool((cap.get("whiteboard") or {}).get("palette_ref"))
        if not (has_slide and has_pal):
            bad.append(st["id"])
    rec("catalog", f"verified_dual complete (n={len(vd)})", not bad, "; ".join(bad) or "all have both sides")

    # no residual 'course' key anywhere in styles
    raw = open("catalog/styles.json").read().lower()
    rec("catalog", "no residual course in styles.json", "course" not in raw,
        "found 'course'" if "course" in raw else "clean")

# ---------- shared: run the shipped SVG validator ----------
def validate_svg(svg_text):
    with tempfile.NamedTemporaryFile("w", suffix=".svg", delete=False, dir=ROOT) as f:
        f.write(svg_text); path = f.name
    try:
        p = subprocess.run(["bash", "scripts/whiteboard-cli.sh", os.path.relpath(path, ROOT)],
                           capture_output=True, text=True, timeout=60)
        return p.returncode == 0, p.stdout + p.stderr
    finally:
        os.unlink(path)

def canonical_svg(colors):
    c = lambda k, d: colors.get(k, d)
    return f'''<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1600 1000">
  <rect x="0" y="0" width="1600" height="1000" style="fill:{c('background','#ffffff')};"/>
  <defs><marker id="a" markerWidth="12" markerHeight="9" refX="10" refY="4.5" orient="auto">
    <polygon points="0 0, 12 4.5, 0 9" style="fill:{c('accent','#333333')};"/></marker></defs>
  <rect x="120" y="120" width="400" height="220" rx="8" style="fill:{c('surface','#eeeeee')}; stroke:{c('primary','#222222')}; stroke-width:3;"/>
  <text x="140" y="200" style="font-size:32px; fill:{c('text','#111111')};">Node A</text>
  <text x="140" y="250" style="font-size:20px; fill:{c('text_muted','#666666')};">caption</text>
  <line x1="520" y1="230" x2="700" y2="230" style="stroke:{c('secondary','#888888')}; stroke-width:3;" marker-end="url(#a)"/>
  <ellipse cx="900" cy="230" rx="120" ry="70" style="fill:{c('primary','#222222')}; stroke:{c('text','#111111')}; stroke-width:3;"/>
  <text x="900" y="238" style="font-size:24px; fill:{c('background','#ffffff')}; text-anchor:middle;">Node B</text>
</svg>'''

# ---------- 2. PALETTE SVG VALIDATION ----------
def check_palettes():
    try:
        palettes = json.load(open("catalog/whiteboard-index.json"))["palettes"]
    except Exception as e:
        rec("palettes", "load", False, str(e)); return
    t0 = time.time(); npass = 0; fails = []
    for p in palettes:
        ok, out = validate_svg(canonical_svg(p.get("colors", {})))
        if ok: npass += 1
        else: fails.append(p["id"])
    dt = time.time() - t0
    rec("palettes", f"{len(palettes)} palettes -> valid SVG",
        npass == len(palettes), f"{npass}/{len(palettes)} pass, {dt:.1f}s" + ("; FAIL: " + ",".join(fails) if fails else ""))

# ---------- 3. VALIDATOR CORRECTNESS ----------
def check_validator():
    good = canonical_svg({"primary": "#2e4a2a", "background": "#efe7d4", "text": "#1a1a17",
                          "surface": "#e6dcc4", "accent": "#d27e96", "secondary": "#3a5a36", "text_muted": "#3a5a36"})
    ok, _ = validate_svg(good)
    rec("validator", "good SVG (marker polygon) passes", ok, "" if ok else "false negative!")

    base = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1600 1000">{}</svg>'
    bads = {
        "<path>": '<path d="M10 10 L20 20" style="stroke:#000;"/>',
        "<linearGradient>": '<defs><linearGradient id="g"/></defs><rect width="10" height="10" style="fill:#000;"/>',
        "<filter>": '<defs><filter id="f"/></defs><rect width="10" height="10" style="fill:#000;"/>',
        "opacity attr": '<rect width="10" height="10" opacity="0.5" style="fill:#000;"/>',
        "rgba color (attr)": '<rect width="10" height="10" fill="rgba(0,0,0,0.5)"/>',
        "rgba color (in style=)": '<rect width="10" height="10" style="fill:rgba(0,0,0,0.5);"/>',
        "opacity (in style=)": '<rect width="10" height="10" style="fill:#000; opacity:0.5;"/>',
        "fill-opacity attr": '<rect width="10" height="10" fill="#000" fill-opacity="0.5"/>',
        "standalone <polygon>": '<polygon points="0 0, 10 0, 5 10" style="fill:#000;"/>',
    }
    for name, frag in bads.items():
        ok, _ = validate_svg(base.format(frag))
        rec("validator", f"rejects {name}", not ok, "false positive (accepted bad)!" if ok else "correctly rejected")

# ---------- 4. TEMPLATE RENDER ----------
def check_render():
    tpls = sorted(glob.glob("templates/slides/gallery/*/template.html"))
    rec("render", "gallery template count == 34", len(tpls) == 34, f"got {len(tpls)}")
    if NO_RENDER or not CHROME:
        rec("render", f"{len(tpls)} gallery templates", True,
            "SKIPPED (--no-render)" if NO_RENDER else "SKIPPED (no Chrome found)")
        return
    t0 = time.time(); npass = 0; fails = []
    outdir = tempfile.mkdtemp()
    for tp in tpls:
        sid = tp.split("/")[-2]
        png = os.path.join(outdir, f"{sid}.png")
        try:
            subprocess.run([CHROME, "--headless", "--disable-gpu", "--no-sandbox",
                            "--hide-scrollbars", "--window-size=1920,1080",
                            "--virtual-time-budget=1500", f"--screenshot={png}",
                            "file://" + os.path.abspath(tp)],
                           capture_output=True, timeout=40)
            if os.path.exists(png) and os.path.getsize(png) > 5000:
                npass += 1
            else:
                fails.append(sid)
        except Exception:
            fails.append(sid)
    dt = time.time() - t0
    rec("render", f"{len(tpls)} gallery templates render non-blank",
        npass == len(tpls), f"{npass}/{len(tpls)} pass, {dt:.0f}s" + ("; FAIL: " + ",".join(fails) if fails else ""))

# ---------- 5. SCRIPTS SMOKE ----------
def check_scripts():
    p = subprocess.run(["bash", "scripts/preflight.sh"], capture_output=True, text=True)
    rec("scripts", "preflight.sh exit 0", p.returncode == 0, "")
    p = subprocess.run(["bash", "scripts/render.sh", "bogus", "x", "y"], capture_output=True, text=True)
    rec("scripts", "render.sh rejects bad mode", p.returncode != 0, f"exit={p.returncode}")
    p = subprocess.run(["bash", "scripts/serve.sh", "bogus"], capture_output=True, text=True)
    rec("scripts", "serve.sh rejects bad mode", p.returncode != 0, f"exit={p.returncode}")
    # course must be gone from active scripts/skill files
    leftovers = []
    for f in ["SKILL.md", "skill.json", "scripts/render.sh", "scripts/serve.sh"]:
        if "course" in open(f).read().lower():
            leftovers.append(f)
    rec("scripts", "no course in active routing files", not leftovers, ",".join(leftovers) or "clean")

# ---------- 6. PPTX BUILD ----------
def check_pptx():
    try:
        import pptx  # noqa
    except ImportError:
        rec("pptx", "python-pptx available", False, "not installed (pip install python-pptx)")
        return
    spec = {"title": "stress", "slides": [
        {"type": "cover", "eyebrow": "TEST", "title": "压测 Cover", "subtitle": "editable check", "footer": "x"},
        {"type": "content", "heading": "Heading 标题", "lead": "lead", "bullets": ["one 一", "two 二"]},
        {"type": "two_col", "heading": "对比", "left_title": "L", "left": ["a"], "right_title": "R", "right": ["b"]},
        {"type": "stat", "heading": "数", "items": [{"value": "10¹²", "label": "x"}, {"value": "9", "label": "y"}]},
        {"type": "statement", "text": "punch line", "sub": "s"},
        {"type": "closing", "title": "end 结束", "subtitle": "s"},
    ]}
    import tempfile, json as _json
    d = tempfile.mkdtemp()
    sp = os.path.join(d, "spec.json"); out = os.path.join(d, "deck.pptx")
    open(sp, "w").write(_json.dumps(spec, ensure_ascii=False))
    p = subprocess.run(["python3", "scripts/build-pptx.py", "--style", "bold-signal",
                        "--spec", sp, "--out", out], capture_output=True, text=True)
    rec("pptx", "build-pptx.py runs", p.returncode == 0, (p.stderr or p.stdout)[-120:] if p.returncode else "ok")
    if p.returncode != 0 or not os.path.exists(out):
        return
    # reopen: must be editable native text (real runs), not images
    from pptx import Presentation
    pr = Presentation(out)
    n = len(pr.slides._sldIdLst)
    rec("pptx", "slide count == 6", n == 6, f"got {n}")
    has_text = any(sh.has_text_frame and sh.text_frame.text.strip()
                   for s in pr.slides for sh in s.shapes)
    rec("pptx", "slides contain editable text", has_text, "real text frames" if has_text else "no text!")
    # 16:9
    ratio_ok = abs((pr.slide_width / pr.slide_height) - 16/9) < 0.01
    rec("pptx", "16:9 stage", ratio_ok, "")

# ---------- run ----------
for fn in (check_catalog, check_palettes, check_validator, check_render, check_pptx, check_scripts):
    try:
        fn()
    except Exception as e:
        rec(fn.__name__, "CRASHED", False, repr(e))

# report
print("\n" + "=" * 72)
print("  sota-present 压测报告 / STRESS TEST REPORT")
print("=" * 72)
cur = None
npass = nfail = 0
for section, name, ok, detail in results:
    if section != cur:
        print(f"\n[{section}]"); cur = section
    flag = "✅" if ok else "❌"
    npass += ok; nfail += not ok
    print(f"  {flag} {name}" + (f"  ·  {detail}" if detail else ""))
print("\n" + "=" * 72)
print(f"  TOTAL: {npass} passed, {nfail} failed")
print("=" * 72)
sys.exit(0 if nfail == 0 else 1)
