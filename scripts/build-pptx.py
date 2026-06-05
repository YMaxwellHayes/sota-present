#!/usr/bin/env python3
"""
build-pptx.py — compile a slide spec + a sota-present design profile into an
EDITABLE .pptx (native text frames + shapes, not images).

The renderer is PROFILE-DRIVEN: it reads a per-style design profile
(catalog/profiles/<id>.json, distilled from the gallery template's design.md)
for the real type scale, colour roles, spacing, chrome and component treatments,
so each style's deck carries its design DNA. Styles without a profile fall back
to a token-derived default profile (so any catalog style still works).

Usage:
  python3 scripts/build-pptx.py --style <style-id> --spec <spec.json> --out out.pptx [--cjk-font "PingFang SC"]

Spec: { "title": "...", "slides": [ {type, ...}, ... ] }
Types: cover | section | agenda | content | two_col | stat | statement | closing | quote
"""
import argparse, json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU_IN = 914400
W, H = 13.333, 7.5
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

def hx(v): return RGBColor.from_string(str(v).lstrip("#").upper())
def fam(tok): return str(tok).split(",")[0].strip().strip('"').strip("'")
def lum(hexv):
    s = str(hexv).lstrip("#"); r, g, b = int(s[0:2],16), int(s[2:4],16), int(s[4:6],16)
    return (0.2126*r + 0.7152*g + 0.0722*b) / 255.0

def upper_ascii(t):  # uppercase Latin only, leave CJK untouched
    return "".join(c.upper() if c.isascii() else c for c in t)

def set_font(run, name=None, ea=None, size=None, bold=None, italic=None, color=None, spacing_pt=None):
    f = run.font
    if size is not None: f.size = Pt(size)
    if bold is not None: f.bold = bold
    if italic is not None: f.italic = italic
    if color is not None: f.color.rgb = hx(color)
    rPr = run._r.get_or_add_rPr()
    if name is not None:
        f.name = name
        el = rPr.find(qn("a:cs")) or rPr.makeelement(qn("a:cs"), {})
        if el.getparent() is None: rPr.append(el)
        el.set("typeface", name)
    if ea is not None:
        el = rPr.find(qn("a:ea")) or rPr.makeelement(qn("a:ea"), {})
        if el.getparent() is None: rPr.append(el)
        el.set("typeface", ea)
    if spacing_pt is not None:
        rPr.set("spc", str(int(spacing_pt * 100)))   # OOXML spc = 1/100 pt

# ---------------- profile loading ----------------
DEFAULT_SCALE = {
    "hero":{"pt":66,"weight":700,"face":"display","tracking":-0.01,"leading":0.98,"transform":"none"},
    "display":{"pt":52,"weight":700,"face":"display","tracking":-0.01,"leading":1.0,"transform":"none"},
    "headline":{"pt":32,"weight":700,"face":"display","tracking":-0.005,"leading":1.04,"transform":"none"},
    "eyebrow":{"pt":13,"weight":700,"face":"mono","tracking":0.16,"leading":1.0,"transform":"upper"},
    "body":{"pt":18,"weight":400,"face":"body","tracking":0,"leading":1.3,"transform":"none"},
    "bullet":{"pt":18,"weight":400,"face":"body","tracking":0,"leading":1.2,"transform":"none"},
    "caption":{"pt":11,"weight":500,"face":"mono","tracking":0.12,"leading":1.0,"transform":"upper"},
    "stat_number":{"pt":62,"weight":700,"face":"display","tracking":-0.02,"leading":0.92,"transform":"none"},
    "stat_unit":{"pt":30,"weight":700,"face":"display","tracking":-0.01,"leading":0.92,"transform":"none"},
    "quote":{"pt":40,"weight":700,"face":"display","tracking":-0.005,"leading":1.1,"transform":"none"},
    "card_title":{"pt":22,"weight":700,"face":"display","tracking":-0.005,"leading":1.04,"transform":"none"},
}
DEFAULT_SPACING = {"margin_x":0.92,"margin_top":0.6,"margin_bottom":0.52,"gap_card":0.3,
                   "gap_bullet":0.2,"card_pad":0.46,"radius_card":0.05,"rule_pt":1.25,"rule_card_pt":2.0}

def token_to_profile(tokens):
    bg = tokens.get("color_bg","#ffffff"); prim = tokens.get("color_primary","#222222")
    acc = tokens.get("color_accent",prim); surf = tokens.get("color_surface",bg)
    txt = tokens.get("color_text","#1a1a1a"); mut = tokens.get("color_text_muted","#888888")
    ground = {"fill":bg,"heading":prim,"body":txt,"muted":mut,"accent":acc,"surface":surf,"border":prim}
    return {
        "id":"_default","ground_default":"bg",
        "grounds":{"bg":dict(ground),"bg_alt":dict(ground)},   # default: no flip
        "fonts":{"display":fam(tokens.get("font_heading","Arial")),"body":fam(tokens.get("font_body","Arial")),
                 "mono":fam(tokens.get("font_mono",tokens.get("font_body","Arial"))),
                 "cjk":None,"cjk_mono":None,"single_weight_display":False},
        "type_scale":DEFAULT_SCALE,"spacing":DEFAULT_SPACING,
        "components":{"footline":False,"pageno":True,"heading_rule":True,"marker":"square",
                      "card_border":False,"card_fill":"surface","monogram":False},
        "depth":{"mode":"flat","shadow":False,"offset_dx":0,"offset_dy":0,"offset_color":"border"},
        "decor":[{"type":"edge_bar","where":"cover","color":"accent"}],
        "rules":{"no_em_dash":True,"lock_accent":True,"off_black_white":True,"mono_upper":True,
                 "body_max_weight":400,"display_force_bold":False},
    }

def load_profile(style_id):
    styles = json.load(open(os.path.join(ROOT,"catalog/styles.json")))["styles"]
    st = next((s for s in styles if s["id"]==style_id), None)
    if not st: sys.exit(f"❌ style not found: {style_id}")
    ref = (st.get("capabilities",{}).get("slides",{}) or {}).get("profile_ref")
    prof = None
    if ref:
        p = os.path.join(ROOT,"catalog/profiles",ref+".json")
        if os.path.exists(p): prof = json.load(open(p))
    if prof is None:
        prof = token_to_profile(st["tokens"]); prof["_name"]=st["name"]+" (token default)"
    else:
        prof["_name"]=st["name"]
    # fill any missing scale roles / spacing from defaults
    for k,v in DEFAULT_SCALE.items(): prof["type_scale"].setdefault(k,v)
    for k,v in DEFAULT_SPACING.items(): prof["spacing"].setdefault(k,v)
    contrast_guard(prof)
    return prof

def contrast_guard(prof):
    """Ensure each ground's text roles separate from its fill; fix buggy/invisible colours."""
    for gname, g in prof["grounds"].items():
        fl = lum(g["fill"]); dark = fl < 0.5
        for role in ("heading","body","muted","accent","border"):
            if role not in g: continue
            if abs(lum(g[role]) - fl) < 0.16:
                if role in ("heading","accent","border"):
                    # try the other brand colour, else neutral
                    alt = g.get("accent") if role!="accent" else g.get("heading")
                    g[role] = alt if alt and abs(lum(alt)-fl)>=0.16 else ("#EDEFF2" if dark else "#1A1A1A")
                else:
                    g[role] = "#D6DEE6" if dark else "#33352F" if not dark else "#D6DEE6"
        # surface should differ slightly from fill (else cards vanish)
        if "surface" in g and abs(lum(g["surface"])-fl) < 0.04:
            g["surface"] = _shift(g["fill"], dark)

def _shift(hexv, dark):
    s=hexv.lstrip("#"); r,gc,b=int(s[0:2],16),int(s[2:4],16),int(s[4:6],16)
    d=18 if dark else -14
    f=lambda x:max(0,min(255,x+d)); return f"#{f(r):02X}{f(gc):02X}{f(b):02X}"

# ---------------- deck ----------------
class Deck:
    def __init__(self, prof, cjk_override):
        self.p = prof; self.ts = prof["type_scale"]; self.sp = prof["spacing"]
        self.comp = prof["components"]; self.depth = prof["depth"]; self.rules = prof["rules"]
        self.fonts = prof["fonts"]; self.cjk = cjk_override
        self.MX = self.sp["margin_x"]; self.CW = W - 2*self.MX
        self.prs = Presentation()
        self.prs.slide_width = Emu(int(W*EMU_IN)); self.prs.slide_height = Emu(int(H*EMU_IN))
        self.blank = self.prs.slide_layouts[6]

    # ground helpers
    def g(self, ground): return self.p["grounds"].get(ground, self.p["grounds"]["bg"])
    def _slide(self, ground):
        s = self.prs.slides.add_slide(self.blank)
        s.background.fill.solid(); s.background.fill.fore_color.rgb = hx(self.g(ground)["fill"])
        return s

    def _face(self, role):
        fr = self.ts[role]["face"]
        latin = self.fonts.get(fr) or self.fonts.get("body")
        ea = self.cjk or (self.fonts.get("cjk_mono") if fr=="mono" else self.fonts.get("cjk"))
        return latin, ea

    def rect(self, s, x, y, w, h, fill, line=None, lw=1.5, rounded=False, radius=0.05):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None: shp.fill.background()
        else: shp.fill.solid(); shp.fill.fore_color.rgb = hx(fill)
        if line: shp.line.color.rgb = hx(line); shp.line.width = Pt(lw)
        else: shp.line.fill.background()
        shp.shadow.inherit = False
        if rounded:
            try: shp.adjustments[0] = radius
            except Exception: pass
        return shp

    def card(self, s, x, y, w, h, ground):
        gg = self.g(ground)
        if self.depth["mode"] == "offset" and (self.depth["offset_dx"] or self.depth["offset_dy"]):
            self.rect(s, x+self.depth["offset_dx"], y+self.depth["offset_dy"], w, h,
                      gg.get(self.depth.get("offset_color","border"), gg["border"]),
                      rounded=True, radius=self.sp["radius_card"])
        fill = gg.get(self.comp.get("card_fill","surface"), gg["surface"])
        line = gg["border"] if self.comp.get("card_border") else None
        self.rect(s, x, y, w, h, fill, line=line, lw=self.sp["rule_card_pt"],
                  rounded=True, radius=self.sp["radius_card"])

    def role_text(self, s, x, y, w, h, text, role, ground, color_role=None,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        rs = self.ts[role]; gg = self.g(ground)
        # default colour per role
        if color_role is None:
            color_role = {"eyebrow":"accent","caption":"muted","body":"body","bullet":"body",
                          "stat_unit":"accent"}.get(role, "heading")
        color = gg.get(color_role, gg["heading"])
        t = str(text)
        if self.rules.get("no_em_dash"): t = t.replace("—"," ").replace("–","-")
        if rs.get("transform")=="upper" and self.rules.get("mono_upper", True): t = upper_ascii(t)
        weight = rs.get("weight",400)
        bold = (weight >= 600) and not (rs["face"]=="display" and self.fonts.get("single_weight_display"))
        if rs["face"]=="body" and weight > self.rules.get("body_max_weight",700): bold=False
        latin, ea = self._face(role)
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
        tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
        for i, line in enumerate(t.split("\n")):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment=align; p.line_spacing=rs.get("leading",1.1)
            r = p.add_run(); r.text = line
            set_font(r, name=latin, ea=ea, size=rs["pt"], bold=bold, italic=rs.get("italic",False),
                     color=color, spacing_pt=rs.get("tracking",0)*rs["pt"])
        return tb

    # chrome
    def pageno(self, s, idx, total, ground):
        if not self.comp.get("pageno"): return
        self.role_text(s, W-2.1, H-0.58, 1.65, 0.4, f"{idx:02d} / {total:02d}", "caption",
                       ground, color_role="muted", align=PP_ALIGN.RIGHT)
    def footline(self, s, title, ground):
        if not self.comp.get("footline"): return
        self.role_text(s, self.MX, H-0.58, self.CW-2.2, 0.4, title, "caption", ground, color_role="muted")

    def heading(self, s, txt, eyebrow, ground):
        gg = self.g(ground); top = self.sp["margin_top"]
        y = top
        if eyebrow:
            self.role_text(s, self.MX+0.02, y, self.CW, 0.4, eyebrow, "eyebrow", ground)
            y += 0.42
        self.rect(s, self.MX, y+0.05, 0.11, 0.62, gg["accent"])
        self.role_text(s, self.MX+0.32, y, self.CW-0.4, 0.95, txt, "headline", ground, anchor=MSO_ANCHOR.MIDDLE)
        ry = y + 0.92
        if self.comp.get("heading_rule"):
            self.rect(s, self.MX, ry, self.CW, self.sp["rule_pt"]/72.0, gg["muted"])
        return ry + 0.28

    def marker(self, s, x, y, ground, size=0.16):
        gg = self.g(ground); kind = self.comp.get("marker","square")
        if kind=="none": return
        if kind=="circle":
            shp=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x),Inches(y),Inches(size),Inches(size))
            shp.fill.solid(); shp.fill.fore_color.rgb=hx(gg["accent"]); shp.line.fill.background(); shp.shadow.inherit=False
        elif kind=="dash":
            self.rect(s, x, y+size*0.4, size*1.3, size*0.22, gg["accent"])
        elif kind=="bar":
            self.rect(s, x, y, size*0.3, size*1.4, gg["accent"])
        else:
            self.rect(s, x, y, size, size, gg["accent"])

    def decor(self, s, where, ground):
        gg = self.g(ground)
        for d in self.p.get("decor", []):
            if d.get("where")!=where: continue
            if d["type"]=="edge_bar":
                self.rect(s, 0, 0, 4.6, 0.14, gg.get(d.get("color","accent"), gg["accent"]))

    # ---------------- renderers ----------------
    def cover(self, d, i, n):
        gd = self.p.get("cover_ground", self.p["ground_default"]); s=self._slide(gd)
        self.decor(s, "cover", gd)
        y = 2.35
        if d.get("eyebrow"): self.role_text(s, self.MX, y, self.CW, 0.5, d["eyebrow"], "eyebrow", gd); y+=0.5
        self.role_text(s, self.MX-0.05, y, self.CW+0.2, 2.6, d["title"], "hero", gd);
        if d.get("subtitle"):
            self.role_text(s, self.MX, 5.25, self.CW-1.0, 1.4, d["subtitle"], "body", gd)
        if d.get("footer"):
            self.role_text(s, self.MX, H-0.72, self.CW, 0.4, d["footer"], "caption", gd, color_role="muted")

    def section(self, d, i, n):
        gd = self.p.get("section_ground", "bg_alt"); s=self._slide(gd); gg=self.g(gd)
        if d.get("num"): self.role_text(s, self.MX-0.05, 0.95, 6, 2.4, d["num"], "hero", gd, color_role="accent")
        self.rect(s, self.MX, 3.4, 0.11, 0.7, gg["accent"])
        self.role_text(s, self.MX+0.34, 3.32, self.CW, 1.1, d["title"], "display", gd, anchor=MSO_ANCHOR.MIDDLE)
        if d.get("subtitle"):
            self.role_text(s, self.MX+0.34, 4.7, self.CW-1.4, 1.2, d["subtitle"], "body", gd, color_role="muted")
        self.pageno(s, i, n, gd)

    def agenda(self, d, i, n):
        gd="bg"; s=self._slide(gd); y=self.heading(s, d["heading"], d.get("kicker"), gd)
        items=d.get("items",[]); y0=max(y,2.2); step=min(0.92,(H-y0-0.65)/max(len(items),1))
        for k,it in enumerate(items):
            self.role_text(s, self.MX+0.02, y0+k*step, 0.7, step, f"{k+1:02d}", "eyebrow", gd, color_role="accent", anchor=MSO_ANCHOR.MIDDLE)
            self.role_text(s, self.MX+0.95, y0+k*step, self.CW-1.0, step, it, "body", gd, anchor=MSO_ANCHOR.MIDDLE)
            if k<len(items)-1: self.rect(s, self.MX+0.95, y0+(k+1)*step-0.02, self.CW-1.0, self.sp["rule_pt"]/72.0, self.g(gd)["muted"])
        self.footline(s, d.get("foot",""), gd); self.pageno(s, i, n, gd)

    def content(self, d, i, n):
        gd="bg"; s=self._slide(gd); y=self.heading(s, d["heading"], d.get("kicker"), gd)
        if d.get("lead"): self.role_text(s, self.MX+0.02, y, self.CW-0.2, 1.0, d["lead"], "body", gd); y+=1.05
        bullets=d.get("bullets",[])
        if bullets:
            avail=H-y-0.6; step=min(0.95,max(0.6,avail/len(bullets)))
            for b in bullets:
                self.marker(s, self.MX+0.04, y+0.16, gd)
                self.role_text(s, self.MX+0.5, y, self.CW-0.6, step, b, "bullet", gd)
                y+=step
        self.footline(s, d.get("foot",""), gd); self.pageno(s, i, n, gd)

    def two_col(self, d, i, n):
        gd="bg"; s=self._slide(gd); top=self.heading(s, d["heading"], d.get("kicker"), gd)
        ct=top+0.05; ch=H-ct-0.65; gap=self.sp["gap_card"]; cw=(self.CW-gap)/2
        for c,(tk,ik) in enumerate([("left_title","left"),("right_title","right")]):
            x=self.MX+c*(cw+gap); self.card(s, x, ct, cw, ch, gd); pad=self.sp["card_pad"]
            self.role_text(s, x+pad, ct+0.34, cw-2*pad, 1.0, d.get(tk,""), "card_title", gd, color_role="accent")
            self.rect(s, x+pad, ct+1.32, cw-2*pad, self.sp["rule_pt"]/72.0, self.g(gd)["muted"])
            iy=ct+1.6; items=d.get(ik,[]); istep=min(0.8,max(0.5,(ct+ch-0.4-iy)/max(len(items),1)))
            for it in items:
                self.marker(s, x+pad, iy+0.1, gd, size=0.13)
                self.role_text(s, x+pad+0.32, iy, cw-2*pad-0.32, istep, it, "bullet", gd)
                iy+=istep
        self.pageno(s, i, n, gd)

    def stat(self, d, i, n):
        gd="bg"; s=self._slide(gd); top=self.heading(s, d.get("heading",""), d.get("kicker"), gd) if d.get("heading") else 2.4
        items=d.get("items",[]); m=max(1,len(items)); gap=0.45; cw=(self.CW-(m-1)*gap)/m; y=max(top,2.5)
        # stat figures are meant to be short; if a value is long (e.g. a name), scale the
        # whole row down so it stays on one line and never overlaps the label below.
        base=self.ts["stat_number"]["pt"]
        longest=max((len(str(it.get("value",""))) for it in items), default=2)
        # rough chars-that-fit at base size in one column; shrink proportionally, clamp
        fit=max(2.0, cw/ (base*0.62/72.0))    # approx Latin glyphs per column width
        size=base if longest<=fit else max(30, base*fit/longest)
        self._stat_pt=size
        for k,it in enumerate(items):
            x=self.MX+k*(cw+gap); self.rect(s, x, y, 0.7, 0.09, self.g(gd)["accent"])
            self._role_text_sized(s, x, y+0.22, cw, size*1.5/72.0+0.1, it.get("value",""), "stat_number", gd, size)
            self.role_text(s, x, y+0.22+size*1.5/72.0+0.18, cw, 1.4, it.get("label",""), "body", gd, color_role="muted")
        self.pageno(s, i, n, gd)

    def _role_text_sized(self, s, x, y, w, h, text, role, ground, size):
        """role_text but with an explicit point size override (for adaptive stat figures)."""
        rs=dict(self.ts[role]); rs["pt"]=size; saved=self.ts[role]; self.ts[role]=rs
        try: self.role_text(s, x, y, w, h, text, role, ground)
        finally: self.ts[role]=saved

    def quote(self, d, i, n):
        gd=self.p.get("section_ground","bg_alt"); s=self._slide(gd); gg=self.g(gd)
        self.rect(s, 0,0,0.14,H, gg["accent"])
        self.role_text(s, self.MX+0.25, 1.7, self.CW-0.4, 3.4, d.get("text",d.get("quote","")), "quote", gd, anchor=MSO_ANCHOR.MIDDLE)
        if d.get("attribution"): self.role_text(s, self.MX+0.25, 5.5, self.CW-1, 0.8, d["attribution"], "eyebrow", gd)
        self.pageno(s, i, n, gd)

    def statement(self, d, i, n):
        gd=self.p.get("statement_ground","bg_alt"); s=self._slide(gd); gg=self.g(gd)
        self.rect(s, 0,0,0.14,H, gg["accent"])
        self.role_text(s, self.MX+0.25, 1.7, self.CW-0.4, 3.4, d["text"], "display", gd, anchor=MSO_ANCHOR.MIDDLE)
        if d.get("sub"): self.role_text(s, self.MX+0.25, 5.45, self.CW-1.3, 1.2, d["sub"], "body", gd, color_role="muted")
        self.pageno(s, i, n, gd)

    def closing(self, d, i, n):
        gd=self.p.get("cover_ground", self.p["ground_default"]); s=self._slide(gd); gg=self.g(gd)
        self.role_text(s, self.MX-0.05, 2.4, self.CW+0.2, 2.6, d["title"], "display", gd)
        if d.get("subtitle"): self.role_text(s, self.MX, 5.1, self.CW-1.0, 1.3, d["subtitle"], "body", gd, color_role="muted")
        self.rect(s, W-4.4, H-0.14, 4.4, 0.14, gg["accent"])

    def render(self, spec):
        fn={"cover":self.cover,"section":self.section,"agenda":self.agenda,"content":self.content,
            "two_col":self.two_col,"stat":self.stat,"statement":self.statement,"closing":self.closing,
            "quote":self.quote}
        title=spec.get("title",""); slides=spec.get("slides",[]); n=len(slides)
        for i,sl in enumerate(slides,1):
            sl.setdefault("foot", title)
            fn.get(sl.get("type","content"), self.content)(sl, i, n)

def main():
    ap=argparse.ArgumentParser()
    ap.add_argument("--style",required=True); ap.add_argument("--spec",required=True)
    ap.add_argument("--out",required=True); ap.add_argument("--cjk-font",default="PingFang SC")
    a=ap.parse_args()
    prof=load_profile(a.style); spec=json.load(open(a.spec))
    os.makedirs(os.path.dirname(os.path.abspath(a.out)),exist_ok=True)
    deck=Deck(prof, a.cjk_font); deck.render(spec); deck.prs.save(a.out)
    print(f"✅ {a.out}  ({len(spec.get('slides',[]))} slides, profile: {prof.get('_name')})")

if __name__=="__main__":
    main()
